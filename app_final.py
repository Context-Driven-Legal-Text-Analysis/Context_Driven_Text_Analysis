from flask import Flask, render_template, request, jsonify
import os
import logging
import PyPDF2
import docx
from pptx import Presentation
import re
import json
import requests
import spacy
import nltk
from utils import count_tokens
from typing import List, Dict, Any
from groq import Groq
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeVectorStore
from rank_bm25 import BM25Okapi
from sentence_transformers import CrossEncoder
import numpy as np
from openai import OpenAI
import chardet
from dotenv import load_dotenv
from deep_translator import GoogleTranslator
from langdetect import detect
import speech_recognition as sr

load_dotenv()

app = Flask(__name__)
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_PATH = os.path.join(BASE_DIR, "uploads")
os.makedirs(UPLOAD_PATH, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_PATH

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

nltk.download('punkt', quiet=True)
nlp = spacy.load("en_core_web_sm")
nlp.max_length = 2000000

docsearch = None
chunks = []


class EmbeddingWrapper:
    def __init__(self, api_key):
        self.api_key = api_key
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://api.deepinfra.com/v1/openai"
        )

    def create_embeddings(self, texts: List[str], batch_size: int = 10) -> List[List[float]]:
        embeddings = []
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            response = self.client.embeddings.create(
                model="BAAI/bge-base-en-v1.5",
                input=batch,
                encoding_format="float"
            )
            embeddings.extend([res.embedding for res in response.data])
        if embeddings:
            print(f"Embedding dimension: {len(embeddings[0])}")
        return embeddings

    def embed_query(self, query: str) -> List[float]:
        response = self.client.embeddings.create(
            model="BAAI/bge-base-en-v1.5",
            input=[query],
            encoding_format="float"
        )
        return response.data[0].embedding


def improved_get_relevant_chunks(user_question: str, docsearch: PineconeVectorStore, chunks: List[str], top_k: int = 3):
    debug_data = {}
    tokenized_query = user_question.split()
    debug_data["query_tokens"] = tokenized_query

    vector_results = docsearch.similarity_search(user_question, k=top_k * 2)
    vector_chunks = [doc.page_content for doc in vector_results]
    debug_data["vector_search_matches"] = vector_chunks

    bm25 = BM25Okapi([doc.split() for doc in chunks])
    bm25_scores = bm25.get_scores(tokenized_query)
    top_bm25_indices = np.argsort(bm25_scores)[-top_k * 2:][::-1]

    bm25_debug = []
    bm25_chunks = []
    for idx in top_bm25_indices:
        bm25_chunks.append(chunks[idx])
        bm25_debug.append({"chunk_preview": chunks[idx][:300], "bm25_score": float(bm25_scores[idx])})
    debug_data["bm25_matches"] = bm25_debug

    combined_chunks = list(dict.fromkeys(vector_chunks + bm25_chunks))
    debug_data["combined_before_rerank"] = combined_chunks

    cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    pairs = [(user_question, chunk) for chunk in combined_chunks]
    cross_encoder_scores = cross_encoder.predict(pairs)

    rerank_debug = []
    for chunk, score in zip(combined_chunks, cross_encoder_scores):
        rerank_debug.append({"chunk_preview": chunk[:300], "cross_score": float(score)})
    debug_data["cross_encoder_scores"] = rerank_debug

    sorted_chunks = [chunk for _, chunk in sorted(zip(cross_encoder_scores, combined_chunks), reverse=True)]
    final_chunks = sorted_chunks[:top_k]
    debug_data["final_selected_chunks"] = final_chunks
    return final_chunks, debug_data


def semantic_chunking(text: str) -> List[str]:
    doc = nlp(text)
    chunks = []
    current_chunk = []
    current_length = 0
    max_chunk_length = 500

    for sent in doc.sents:
        sent_length = len(sent)
        if current_length + sent_length > max_chunk_length and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0
        current_chunk.append(sent.text)
        current_length += sent_length

    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return chunks


def populate_pinecone_index(file_path, embedding_wrapper, pc, index_name):
    try:
        text = extract_text_from_files(file_path)
        if not text:
            raise ValueError(f"No text could be extracted from: {file_path}")

        chunks = semantic_chunking(text)
        embeddings = embedding_wrapper.create_embeddings(chunks)

        if len(embeddings[0]) != 768:
            raise ValueError(f"Embedding dimension {len(embeddings[0])} does not match index dimension 768")

        data_to_upsert = [
            (f"chunk_{i}", embedding, {"text": chunk})
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
        ]

        index = pc.Index(index_name)
        index.upsert(vectors=data_to_upsert)
        logger.info(f"Upserted {len(data_to_upsert)} vectors to Pinecone index")

        vector_store = PineconeVectorStore(index, embedding_wrapper, "text")
        return vector_store, chunks
    except Exception as e:
        logger.error(f"Error populating Pinecone index: {str(e)}")
        raise


def check_index_populated(index):
    stats = index.describe_index_stats()
    return stats.total_vector_count > 0


def clean_text(text):
    text = re.sub(r'\s+', ' ', text).strip()
    text = ''.join(char for char in text if char.isprintable() or char.isspace())
    return text


def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {str(e)}")
    return clean_text(text)


def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {str(e)}")
        return ""


def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {str(e)}")
        return ""


def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from TXT {txt_path}: {str(e)}")
        return ""


def extract_text_from_files(path):
    all_text = ""
    if os.path.isdir(path):
        for filename in os.listdir(path):
            file_path = os.path.join(path, filename)
            logger.info(f"Processing file: {filename}")
            if filename.lower().endswith('.pdf'):
                text = extract_text_from_pdf(file_path)
            elif filename.lower().endswith('.docx'):
                text = extract_text_from_docx(file_path)
            elif filename.lower().endswith('.pptx'):
                text = extract_text_from_pptx(file_path)
            elif filename.lower().endswith('.txt'):
                text = extract_text_from_txt(file_path)
            else:
                logger.warning(f"Unsupported file type: {filename}")
                continue
            if text:
                all_text += f"--- Start of {filename} ---\n{text}\n--- End of {filename} ---\n\n"
    elif os.path.isfile(path):
        filename = os.path.basename(path)
        if filename.lower().endswith('.pdf'):
            text = extract_text_from_pdf(path)
        elif filename.lower().endswith('.docx'):
            text = extract_text_from_docx(path)
        elif filename.lower().endswith('.pptx'):
            text = extract_text_from_pptx(path)
        elif filename.lower().endswith('.txt'):
            text = extract_text_from_txt(path)
        else:
            text = None
        if text:
            all_text += f"--- Start of {filename} ---\n{text}\n--- End of {filename} ---\n\n"
    else:
        logger.error(f"Invalid path: {path}")
    return all_text


def rewrite_query(client, model, original_query, rewrite_conversation_history):
    system_prompt = '''You are a document analysis assistant. Rewrite the user query to improve document retrieval.'''
    messages = [
        {"role": "system", "content": system_prompt},
        *rewrite_conversation_history,
        {"role": "user", "content": f"Original query: {original_query}\n\nRewritten query:"}
    ]
    response = client.chat.completions.create(model=model, messages=messages, max_tokens=200, temperature=0.3)
    rewritten_query = response.choices[0].message.content
    rewrite_conversation_history.append({"role": "user", "content": original_query})
    rewrite_conversation_history.append({"role": "assistant", "content": rewritten_query})
    if len(rewrite_conversation_history) > 20:
        rewrite_conversation_history = rewrite_conversation_history[-20:]
    return rewritten_query, rewrite_conversation_history


def chatbot_response(client, model, user_question, relevant_chunks, conversation_history):
    system_prompt = '''You are an AI assistant. Answer ONLY using provided document chunks.
If answer not in chunks, say: "Not found in document".'''
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"Question: {user_question}\n\nContext from document:\n{relevant_chunks}\n\nAnswer:"}
    ]
    response = client.chat.completions.create(model=model, messages=messages, max_tokens=2000, temperature=0.2)
    answer = response.choices[0].message.content
    conversation_history.append({"role": "user", "content": user_question})
    conversation_history.append({"role": "assistant", "content": answer})
    if len(conversation_history) > 2:
        conversation_history = conversation_history[-2:]
    return answer


def translate_text(text: str, source_lang: str, target_lang: str) -> str:
    """Translate text between any two languages. Returns original if langs match or on error."""
    if source_lang == target_lang:
        return text
    try:
        return GoogleTranslator(source=source_lang, target=target_lang).translate(text)
    except Exception as e:
        logger.warning(f"Translation error ({source_lang} → {target_lang}): {e}")
        return text


def translate_to_english(text: str):
    """Detect language and translate to English. Returns (english_text, detected_lang)."""
    try:
        detected_lang = detect(text)
        english_text = translate_text(text, detected_lang, 'en') if detected_lang != 'en' else text
        return english_text, detected_lang
    except Exception as e:
        logger.warning(f"translate_to_english error: {e}")
        return text, 'en'


@app.route('/')
def index():
    return render_template('index_final.html')


@app.route('/start_recording', methods=['POST'])
def start_recording():
    try:
        data = request.get_json() or {}
        # FIX: read the user's chosen output language from the dropdown
        target_lang = data.get('language', 'en')

        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.listen(source, timeout=5)

        # Google STT: returns text in the spoken language
        spoken_text = recognizer.recognize_google(audio)

        # Detect what language was spoken
        detected_lang = detect(spoken_text)

        # Always translate spoken text → English for the chat input box
        english_text = translate_text(spoken_text, detected_lang, 'en')

        # FIX: also produce a version in the user's chosen display language
        display_text = translate_text(english_text, 'en', target_lang) if target_lang != 'en' else english_text

        return jsonify({
            'success': True,
            'text': english_text,        # sent to /chat for retrieval (always English)
            'display_text': display_text, # shown in the input box (user's chosen language)
            'detected_lang': detected_lang
        })

    except sr.UnknownValueError:
        return jsonify({'error': 'Could not understand audio'}), 400
    except sr.RequestError as e:
        return jsonify({'error': f'Speech service error: {str(e)}'}), 500
    except Exception as e:
        logger.error(f"/start_recording error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/stop_recording', methods=['POST'])
def stop_recording():
    return jsonify({'success': True})


@app.route('/upload', methods=['POST'])
def upload_file():
    global docsearch, chunks
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(file_path)
    text = extract_text_from_files(file_path)
    emb = EmbeddingWrapper(api_key=os.getenv("DEEPINFRA_API_KEY"))
    embeddings = emb.create_embeddings([text])

    if not check_index_populated(index):
        docsearch, chunks = populate_pinecone_index(file_path, emb, pc, pinecone_index_name)
    else:
        docsearch = PineconeVectorStore(index, emb, "text")
        chunks = semantic_chunking(text)

    return jsonify({"message": "File uploaded and processed successfully"}), 200


@app.route('/chat', methods=['POST'])
def chat():
    global docsearch, chunks

    if docsearch is None or not chunks:
        return jsonify({"error": "No document uploaded yet. Please upload a file first."}), 400

    body = request.json or {}
    user_message = body.get('message', '')

    # FIX: read the language the user explicitly chose from the dropdown
    # Fall back to auto-detect only if not provided
    target_lang = body.get('target_lang', '').strip()

    # Translate the incoming question to English for retrieval & answering
    user_question_en, detected_lang = translate_to_english(user_message)

    # If the user didn't pass a target_lang, reply in whatever language they wrote in
    if not target_lang:
        target_lang = detected_lang

    # Rewrite query for better retrieval
    rewritten_query, _ = rewrite_query(client, model, user_question_en, [])

    # Retrieve relevant chunks
    relevant_chunks, debug_data = improved_get_relevant_chunks(rewritten_query, docsearch, chunks)

    # Generate English response
    response_en = chatbot_response(client, model, user_question_en, relevant_chunks, [])

    # FIX: translate the response to the user's CHOSEN language, not a random detected one
    response_final = translate_text(response_en, 'en', target_lang)

    return jsonify({
        "response": response_final,
        "target_lang": target_lang,
        "detected_lang": detected_lang,
        "rewritten_query": rewritten_query,
        "retrieved_chunks": relevant_chunks,
        "retrieval_debug": debug_data
    })


if __name__ == "__main__":
    model = os.getenv("MODEL")
    groq_api_key = os.getenv("GROQ_API_KEY")
    pinecone_api_key = os.getenv("PINECONE_API_KEY")
    pinecone_index_name = os.getenv("INDEX_NAME")
    deepinfra_api_key = os.getenv("DEEPINFRA_API_KEY")

    client = Groq(api_key=groq_api_key)
    pc = Pinecone(api_key=pinecone_api_key, request_timeout=60)
    embedding_wrapper = EmbeddingWrapper(deepinfra_api_key)

    if pinecone_index_name not in pc.list_indexes().names():
        pc.create_index(
            pinecone_index_name,
            dimension=768,
            metric='cosine',
            spec=ServerlessSpec(cloud='aws', region='us-east-1')
        )
        logger.info(f"Created new Pinecone index: {pinecone_index_name}")

    index = pc.Index(pinecone_index_name)
    app.run(debug=True)